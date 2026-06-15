#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
读取 JSONL 文件（每行一个 json，含 text / file / slide 键），生成 pptx。

规则：
- 每行 json 生成一页幻灯片。
- 按 file 键的值分组：相同 file 值的所有行写进同一个 pptx，
  页面顺序按 slide 键（slide_0, slide_1, ...）的数字排序。
- 每页的 text 按 512 个字符切分：总长度 <= 512 不切；> 512 则切成多块，
  每块放进一个独立文本框，从上到下依次排列，尽量不超出页面。
- 若存在模板 template.pptx，则参考其背景：为每页克隆模板页的背景
  （<p:bg> 填充、背景图片及图片关系、母版/版式背景），去掉模板里原有
  的占位符/空文本框，再叠加我们自己的文本框。
"""

import argparse
import copy
import glob
import json
import os
import re
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

CHUNK_SIZE = 512

# 没有模板时使用的 16:9 宽屏尺寸
DEFAULT_SLIDE_W = Inches(13.333)
DEFAULT_SLIDE_H = Inches(7.5)

# 内容区边距
MARGIN_X = Inches(0.5)
MARGIN_Y = Inches(0.4)

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# --------------------------------------------------------------------------- #
# 文本切分与排版
# --------------------------------------------------------------------------- #
def split_text(text, size=CHUNK_SIZE):
    """按字符长度切分；<= size 则整体返回一块。"""
    if text is None:
        text = ""
    if len(text) <= size:
        return [text]
    return [text[i:i + size] for i in range(0, len(text), size)]


def pick_font_size(n_chunks):
    """块越多字号越小，尽量塞进一页。"""
    if n_chunks <= 1:
        return 16
    if n_chunks <= 2:
        return 14
    if n_chunks <= 4:
        return 12
    if n_chunks <= 6:
        return 10
    if n_chunks <= 9:
        return 9
    return 8


def fill_text(slide, text, slide_w, slide_h):
    """把 text 切块后逐块放入文本框，自上而下排列，尽量不超出页面。"""
    usable_w = slide_w - 2 * MARGIN_X
    usable_h = slide_h - 2 * MARGIN_Y

    chunks = split_text(text, CHUNK_SIZE)
    n = len(chunks)
    font_size = pick_font_size(n)

    gap = Inches(0.1)
    total_gap = gap * (n - 1) if n > 1 else 0
    box_h = int((usable_h - total_gap) / n)

    top = MARGIN_Y
    for chunk in chunks:
        tb = slide.shapes.add_textbox(MARGIN_X, top, usable_w, box_h)
        tf = tb.text_frame
        tf.word_wrap = True

        # chunk 内部保留原有换行
        first = True
        for line in chunk.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)

        top = top + box_h + gap


# --------------------------------------------------------------------------- #
# 模板背景克隆
# --------------------------------------------------------------------------- #
def _is_placeholder_or_textbox(sp_el):
    """判断一个 <p:sp> 是否为占位符或文本框（这些是要丢弃的“空文本框”）。"""
    nv = sp_el.find(qn("p:nvSpPr"))
    if nv is None:
        return False
    cNvSpPr = nv.find(qn("p:cNvSpPr"))
    if cNvSpPr is not None and cNvSpPr.get("txBox") == "1":
        return True
    nvPr = nv.find(qn("p:nvPr"))
    if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
        return True
    return False


def _remap_relationships(new_part, src_part, elements):
    """把复制过来的元素里引用的 rId（图片等）重新指向新页的关系。"""
    rid_map = {}
    rel_attr_prefix = "{%s}" % R_NS
    for root in elements:
        for el in root.iter():
            for key, val in list(el.attrib.items()):
                if not key.startswith(rel_attr_prefix):
                    continue
                old_rid = val
                if old_rid not in rid_map:
                    rel = src_part.rels[old_rid]
                    if rel.is_external:
                        new_rid = new_part.relate_to(
                            rel.target_ref, rel.reltype, is_external=True)
                    else:
                        new_rid = new_part.relate_to(rel.target_part, rel.reltype)
                    rid_map[old_rid] = new_rid
                el.set(key, rid_map[old_rid])


def make_template_slide(prs, template_slide):
    """基于模板页新建一页：复制背景，去掉占位符/文本框，返回新页。"""
    new_slide = prs.slides.add_slide(template_slide.slide_layout)

    # 移除从版式继承来的占位符（我们要自己加文本框）
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    src_cSld = template_slide._element.find(qn("p:cSld"))
    new_cSld = new_slide._element.find(qn("p:cSld"))
    copied = []

    # 1) 复制幻灯片级背景 <p:bg>（必须位于 spTree 之前）
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        new_bg = copy.deepcopy(src_bg)
        new_cSld.insert(0, new_bg)
        copied.append(new_bg)

    # 2) 复制形状：丢弃占位符/文本框，保留背景图片与装饰形状
    spTree = new_slide.shapes._spTree
    for shp in template_slide.shapes:
        el = shp._element
        if el.tag == qn("p:sp") and _is_placeholder_or_textbox(el):
            continue
        new_el = copy.deepcopy(el)
        spTree.append(new_el)
        copied.append(new_el)

    # 3) 修复复制元素里的图片等关系引用
    _remap_relationships(new_slide.part, template_slide.part, copied)

    return new_slide


def _remove_slides(prs, slides):
    """从演示文稿中删除指定的若干页（用于清掉模板自带的空页）。"""
    target_parts = {s.part for s in slides}
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        if prs.part.related_part(rId) in target_parts:
            sldIdLst.remove(sldId)
            prs.part.drop_rel(rId)


# --------------------------------------------------------------------------- #
# 输入读取
# --------------------------------------------------------------------------- #
def slide_sort_key(slide_val):
    """从 slide_0 / slide_12 之类的值里取出数字用于排序。"""
    if slide_val is None:
        return (1, 0, "")
    m = re.search(r"(\d+)", str(slide_val))
    if m:
        return (0, int(m.group(1)), "")
    return (1, 0, str(slide_val))


def load_records(input_paths):
    """从多个 json/jsonl 文件读出所有记录，逐行解析。"""
    records = []
    for path in input_paths:
        with open(path, "r", encoding="utf-8") as f:
            for lineno, raw in enumerate(f, 1):
                raw = raw.strip()
                if not raw:
                    continue
                try:
                    obj = json.loads(raw)
                except json.JSONDecodeError as e:
                    print(f"[警告] {path}:{lineno} JSON 解析失败，跳过：{e}",
                          file=sys.stderr)
                    continue
                records.append(obj)
    return records


def safe_filename(name):
    """保证 file 值能作为文件名使用，并以 .pptx 结尾。"""
    name = str(name).strip() or "output"
    name = re.sub(r"[\\/:*?\"<>|]", "_", name)
    if not name.lower().endswith(".pptx"):
        name += ".pptx"
    return name


# --------------------------------------------------------------------------- #
# 主流程
# --------------------------------------------------------------------------- #
def build_presentation(recs, template_path):
    """为一个 file 分组构建演示文稿。"""
    if template_path:
        prs = Presentation(template_path)
        original_slides = list(prs.slides)
        template_slide = original_slides[0]
        slide_w, slide_h = prs.slide_width, prs.slide_height

        for rec in recs:
            slide = make_template_slide(prs, template_slide)
            fill_text(slide, rec.get("text", ""), slide_w, slide_h)

        # 清掉模板自带的空页
        _remove_slides(prs, original_slides)
    else:
        prs = Presentation()
        prs.slide_width = DEFAULT_SLIDE_W
        prs.slide_height = DEFAULT_SLIDE_H
        for rec in recs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
            fill_text(slide, rec.get("text", ""),
                      prs.slide_width, prs.slide_height)
    return prs


def main():
    parser = argparse.ArgumentParser(description="按 JSONL 生成 pptx")
    parser.add_argument("-i", "--input-dir", default="input",
                        help="输入目录，扫描其中的 .json/.jsonl（默认 input）")
    parser.add_argument("-o", "--output-dir", default="output",
                        help="输出目录（默认 output）")
    parser.add_argument("-t", "--template", default="template.pptx",
                        help="模板 pptx，参考其背景（默认 template.pptx，"
                             "不存在则用空白主题）")
    args = parser.parse_args()

    in_dir = args.input_dir
    out_dir = args.output_dir
    os.makedirs(out_dir, exist_ok=True)

    template_path = args.template if os.path.isfile(args.template) else None
    if template_path:
        print(f"[信息] 使用模板：{template_path}")
    else:
        print("[信息] 未找到模板，使用空白主题")

    input_paths = sorted(
        glob.glob(os.path.join(in_dir, "*.json")) +
        glob.glob(os.path.join(in_dir, "*.jsonl"))
    )
    if not input_paths:
        print(f"[错误] 在 {in_dir} 下没有找到 .json/.jsonl 文件", file=sys.stderr)
        sys.exit(1)

    records = load_records(input_paths)
    if not records:
        print("[错误] 没有可用记录", file=sys.stderr)
        sys.exit(1)

    # 按 file 值分组（保持首次出现的顺序）
    groups = {}
    order = []
    for rec in records:
        file_val = rec.get("file")
        if file_val is None:
            file_val = "output"
        if file_val not in groups:
            groups[file_val] = []
            order.append(file_val)
        groups[file_val].append(rec)

    for file_val in order:
        recs = groups[file_val]
        recs.sort(key=lambda r: slide_sort_key(r.get("slide")))

        prs = build_presentation(recs, template_path)

        out_path = os.path.join(out_dir, safe_filename(file_val))
        prs.save(out_path)
        print(f"[完成] {out_path}  ({len(recs)} 页)")


if __name__ == "__main__":
    main()
